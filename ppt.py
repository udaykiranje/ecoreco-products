from pptx import Presentation
from pptx.util import Inches, Pt

def create_viva_presentation():
    prs = Presentation()

    # 1. Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "EcoReco: Machine Learning Product Recommendation"
    slide.placeholders[1].text = "Uday Kiran Jella | ID: 23089394\nSupervisor: Sydney Ezika"

    # 2. Data Cleaning & Pre-processing
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Data Pre-processing Steps"
    tf = slide.placeholders[1].text_frame
    tf.text = "• Parsed 'price' by removing '$' and converting to float."
    tf.add_paragraph().text = "• Handled missing 'material' and 'description' values with empty strings."
    tf.add_paragraph().text = "• Combined Title, Material, and Description for semantic analysis."
    tf.add_paragraph().text = "• Scaled numeric features (Rating, Reviews, Price) using MinMaxScaler."

    # 3. Model Implementation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Model: TF-IDF & Hybrid Ranking"
    tf = slide.placeholders[1].text_frame
    tf.text = "• TF-IDF Vectorization: Captured unigrams and bigrams for text weights."
    tf.add_paragraph().text = "• Similarity: Cosine Similarity matrix used for item-to-item distance."
    tf.add_paragraph().text = "• Formula: α(Sim) + β(Green) + γ(Rating) - δ(Price)."
    tf.add_paragraph().text = "• Green Score: Custom keyword matching (e.g., bamboo, organic, recycled)."

    # 4. Results & Performance
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Performance Metrics"
    tf = slide.placeholders[1].text_frame
    tf.text = "• Precision@5: 0.123 | Recall@5: 0.326"
    tf.add_paragraph().text = "• F1-score: 0.178 | MAP@5: 0.333"
    tf.add_paragraph().text = "• Mean Precision for test queries (jute, burlap, cotton): 0.20."

    # 5. Conclusion & Future Work
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Conclusion & Improvements"
    tf = slide.placeholders[1].text_frame
    tf.text = "• Achieved: Interpretable system balancing personalization and ethics."
    tf.add_paragraph().text = "• Next Step: Integrate real-time user interaction data (CF)."
    tf.add_paragraph().text = "• Future: Use CNNs for visual similarity and carbon footprint tracking."

    prs.save('EcoReco_Presentation.pptx')
    print("Document 'EcoReco_Presentation.pptx' generated successfully!")

create_viva_presentation()